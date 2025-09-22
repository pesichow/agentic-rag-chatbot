# document_parser.py
import PyPDF2
from pptx import Presentation
import pandas as pd
from docx import Document
import re
from typing import List, Dict, Any

class DocumentParser:
    def parse_document(self, file_path: str, file_type: str) -> List[Dict[str, Any]]:
        if file_type == "pdf":
            return self._parse_pdf(file_path)
        elif file_type == "pptx":
            return self._parse_pptx(file_path)
        elif file_type == "csv":
            return self._parse_csv(file_path)
        elif file_type == "docx":
            return self._parse_docx(file_path)
        elif file_type in ["txt", "md"]:
            return self._parse_text(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    def _parse_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        chunks = []
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page_num, page in enumerate(reader.pages):
                text = page.extract_text()
                if text.strip():
                    chunks.append({
                        "content": text,
                        "metadata": {"page": page_num + 1, "source": file_path}
                    })
        return chunks
    
    def _parse_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        chunks = []
        presentation = Presentation(file_path)
        for slide_num, slide in enumerate(presentation.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                chunks.append({
                    "content": "\n".join(slide_text),
                    "metadata": {"slide": slide_num + 1, "source": file_path}
                })
        return chunks
    
    def _parse_csv(self, file_path: str) -> List[Dict[str, Any]]:
        chunks = []
        df = pd.read_csv(file_path)
        # Convert DataFrame to text chunks
        text = df.to_string()
        # Split into reasonable chunks
        lines = text.split('\n')
        chunk_size = 20  # lines per chunk
        for i in range(0, len(lines), chunk_size):
            chunk_content = '\n'.join(lines[i:i+chunk_size])
            chunks.append({
                "content": chunk_content,
                "metadata": {"chunk": i//chunk_size + 1, "source": file_path}
            })
        return chunks
    
    def _parse_docx(self, file_path: str) -> List[Dict[str, Any]]:
        chunks = []
        doc = Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
        # Split into paragraphs as chunks
        paragraphs = [p for p in text.split('\n') if p.strip()]
        for i, para in enumerate(paragraphs):
            chunks.append({
                "content": para,
                "metadata": {"paragraph": i + 1, "source": file_path}
            })
        return chunks
    
    def _parse_text(self, file_path: str) -> List[Dict[str, Any]]:
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        # Split into paragraphs
        paragraphs = [p for p in re.split(r'\n\s*\n', content) if p.strip()]
        chunks = []
        for i, para in enumerate(paragraphs):
            chunks.append({
                "content": para,
                "metadata": {"paragraph": i + 1, "source": file_path}
            })
        return chunks